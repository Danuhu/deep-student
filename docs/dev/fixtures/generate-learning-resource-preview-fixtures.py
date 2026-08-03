#!/usr/bin/env python3
from __future__ import annotations

import csv
import json
import math
import shutil
import struct
import wave
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from pypdf import PdfReader, PdfWriter
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[3]
OUT = ROOT / "output" / "learning-resource-preview-fixtures"
TMP = OUT / "_source"
MARKER = "DSTU-PREVIEW-20260720"


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            try:
                return ImageFont.truetype(candidate, size=size, index=0)
            except (OSError, ValueError):
                continue
    return ImageFont.load_default()


def make_illustration(path: Path) -> None:
    image = Image.new("RGB", (1200, 675), "#F4F1E8")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((60, 60, 1140, 615), radius=40, fill="#173F5F")
    draw.ellipse((105, 120, 405, 420), fill="#F6D55C")
    draw.polygon([(300, 500), (570, 170), (840, 500)], fill="#3CAEA3")
    draw.rounded_rectangle((720, 120, 1060, 380), radius=24, fill="#ED553B")
    draw.text((105, 520), "Deep Student 学习资源预览", fill="white", font=font(40, True))
    image.save(path, quality=92)


def base_document(title: str) -> Document:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.5)
    title_paragraph = doc.add_heading(title, 0)
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = doc.add_paragraph(f"预览校验标识：{MARKER}")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    return doc


def make_rich_docx(path: Path, illustration: Path) -> None:
    doc = base_document("跨学科学习笔记：气候、数据与决策")
    doc.add_heading("1. 学习目标", level=1)
    doc.add_paragraph("本资料用于验证中文、English、数学符号（α β γ）、标点和 emoji 🌱 的混排预览。")
    objectives = [
        "解释温室效应的基本机制。",
        "读取并比较年度观测数据。",
        "用证据区分相关性与因果关系。",
    ]
    for item in objectives:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_picture(str(illustration), width=Inches(6.6))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_heading("2. 数据表", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.style = "Light Shading Accent 1"
    for cell, value in zip(table.rows[0].cells, ["年份", "温度异常 (°C)", "CO₂ (ppm)", "观察"]):
        cell.text = value
    for row in [
        ("1990", "0.45", "354", "基线期"),
        ("2010", "0.72", "390", "持续上升"),
        ("2025", "1.18", "425", "需进一步行动"),
    ]:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            cell.text = value
    doc.add_heading("3. 关键公式与引用", level=1)
    doc.add_paragraph("样例公式：E = mc²；函数：f(x) = x² + 2x + 1；化学式：H₂O、CO₂。")
    quote = doc.add_paragraph("“学习的价值在于把信息转化为可检验的理解。”")
    quote.style = "Quote"
    doc.add_page_break()
    doc.add_heading("4. 复习题", level=1)
    for index, question in enumerate(
        [
            "为什么只看单一年份的数据可能产生误导？",
            "表格中的哪个变化最值得进一步调查？",
            "请用 100 字概括本资料的中心观点。",
        ],
        1,
    ):
        doc.add_paragraph(f"{index}. {question}")
        doc.add_paragraph("答题区：" + "_" * 72)
    doc.save(path)


def make_table_docx(path: Path) -> None:
    doc = base_document("横向长表格与分页测试")
    section = doc.sections[0]
    section.orientation = WD_ORIENT.LANDSCAPE
    section.page_width, section.page_height = section.page_height, section.page_width
    doc.add_paragraph("用于验证宽表格、重复表头、分页、数字和长文本是否会被裁切。")
    table = doc.add_table(rows=1, cols=7)
    table.style = "Table Grid"
    headers = ["序号", "章节", "关键词", "掌握度", "计划分钟", "完成日期", "备注"]
    for cell, value in zip(table.rows[0].cells, headers):
        cell.text = value
    for i in range(1, 46):
        values = [
            str(i),
            f"第 {(i - 1) // 5 + 1} 单元",
            f"概念-{i:02d}",
            f"{55 + i % 41}%",
            str(15 + i % 6 * 5),
            f"2026-07-{(i - 1) % 28 + 1:02d}",
            "长文本备注：复习定义、例题与错题，并记录尚未解决的问题。",
        ]
        cells = table.add_row().cells
        for cell, value in zip(cells, values):
            cell.text = value
    doc.save(path)


def make_media_docx(path: Path, illustration: Path) -> None:
    doc = base_document("图文、链接与特殊字符测试")
    doc.add_heading("视觉资料", level=1)
    doc.add_picture(str(illustration), width=Inches(6.9))
    doc.add_paragraph("图片下方说明：色块边缘应清晰，比例应保持 16:9。")
    doc.add_heading("多语言段落", level=1)
    doc.add_paragraph("简体中文｜繁體中文｜日本語｜한국어｜Español｜Français｜العربية")
    doc.add_heading("字符边界", level=1)
    doc.add_paragraph("括号（）【】｛｝；引号“”‘’；货币 ¥ $ €；箭头 ← ↑ → ↓；勾选 ✓。")
    doc.add_paragraph(f"尾部校验标识：{MARKER}-DOCX-MEDIA")
    doc.save(path)


def make_empty_docx(path: Path) -> None:
    doc = Document()
    doc.add_paragraph("")
    doc.save(path)


def register_pdf_font() -> str:
    name = "STSong-Light"
    if name not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(UnicodeCIDFont(name))
    return name


def make_text_pdf(path: Path) -> None:
    pdf_font = register_pdf_font()
    page = canvas.Canvas(str(path), pagesize=A4, pageCompression=1)
    width, height = A4
    for page_no in range(1, 4):
        page.setFillColor(colors.HexColor("#173F5F"))
        page.rect(0, height - 55 * mm, width, 55 * mm, fill=1, stroke=0)
        page.setFillColor(colors.white)
        page.setFont(pdf_font, 24)
        page.drawString(22 * mm, height - 28 * mm, f"文本型 PDF 学习讲义 - 第 {page_no} 页")
        page.setFont(pdf_font, 11)
        page.drawString(22 * mm, height - 42 * mm, f"可选择文字校验：{MARKER}-TEXT-PDF-{page_no}")
        page.setFillColor(colors.HexColor("#1F2933"))
        page.setFont(pdf_font, 15)
        page.drawString(22 * mm, height - 75 * mm, f"{page_no}. 主动学习与证据")
        page.setFont(pdf_font, 11)
        lines = [
            "本页用于验证 PDF 文字层、中文字体、分页、缩放和文本选择。",
            "学习不是重复阅读，而是主动提取、解释、比较并修正理解。",
            "示例数据：数学 45 分钟，语言 30 分钟，科学 60 分钟。",
            "关键问题：哪些证据支持结论？还缺少哪些信息？",
            "English marker: selectable text and page navigation should work.",
        ]
        y = height - 90 * mm
        for index, line in enumerate(lines, 1):
            page.drawString(28 * mm, y, f"{index}. {line}")
            y -= 13 * mm
        page.setStrokeColor(colors.HexColor("#3CAEA3"))
        page.setLineWidth(2)
        page.line(22 * mm, 45 * mm, width - 22 * mm, 45 * mm)
        page.setFillColor(colors.HexColor("#5B6670"))
        page.setFont(pdf_font, 9)
        page.drawRightString(width - 22 * mm, 32 * mm, f"{page_no} / 3")
        page.showPage()
    page.save()


def make_table_pdf(path: Path) -> None:
    pdf_font = register_pdf_font()
    style = ParagraphStyle("body", fontName=pdf_font, fontSize=8, leading=10, textColor=colors.HexColor("#1F2933"))
    title_style = ParagraphStyle("title", fontName=pdf_font, fontSize=18, leading=22, textColor=colors.HexColor("#173F5F"))
    document = SimpleDocTemplate(
        str(path),
        pagesize=landscape(A4),
        leftMargin=14 * mm,
        rightMargin=14 * mm,
        topMargin=12 * mm,
        bottomMargin=12 * mm,
    )
    data: list[list[Paragraph]] = [
        [Paragraph(value, style) for value in ["序号", "章节", "关键词", "掌握度", "分钟", "日期", "备注"]]
    ]
    for i in range(1, 46):
        values = [
            str(i),
            f"第 {(i - 1) // 5 + 1} 单元",
            f"概念-{i:02d}",
            f"{55 + i % 41}%",
            str(15 + i % 6 * 5),
            f"2026-07-{(i - 1) % 28 + 1:02d}",
            "复习定义、例题与错题，记录尚未解决的问题。",
        ]
        data.append([Paragraph(value, style) for value in values])
    table = Table(data, repeatRows=1, colWidths=[15 * mm, 28 * mm, 28 * mm, 20 * mm, 18 * mm, 28 * mm, 98 * mm])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#173F5F")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#91A3B0")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#EEF7F6")]),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    document.build(
        [
            Paragraph("横向宽表格 PDF 预览测试", title_style),
            Spacer(1, 3 * mm),
            Paragraph(f"重复表头、跨页与缩放校验：{MARKER}-TABLE-PDF", style),
            Spacer(1, 4 * mm),
            table,
        ]
    )


def make_scanned_pdf(path: Path) -> None:
    pages: list[Image.Image] = []
    for page_no in range(1, 3):
        image = Image.new("RGB", (1240, 1754), "#FCFBF7")
        draw = ImageDraw.Draw(image)
        draw.rectangle((65, 65, 1175, 1689), outline="#6A6258", width=4)
        draw.text((110, 120), f"扫描讲义样本 - 第 {page_no} 页", fill="#24211E", font=font(52, True))
        draw.text((110, 220), f"图像型 PDF，无文本层｜{MARKER}-SCAN-{page_no}", fill="#403A34", font=font(30))
        y = 330
        for i in range(12):
            draw.text(
                (125, y),
                f"{i + 1:02d}  这是模拟扫描内容，用于检查缩放、翻页与清晰度。",
                fill="#292622",
                font=font(28),
            )
            y += 88
        draw.line((105, 1450, 1135, 1450), fill="#8C8175", width=3)
        draw.text((110, 1510), "手写批注：重点复习第 3、7、10 条。✓", fill="#A23E2C", font=font(34))
        pages.append(image)
    pages[0].save(path, "PDF", resolution=150, save_all=True, append_images=pages[1:])


def make_encrypted_pdf(source: Path, target: Path) -> None:
    reader = PdfReader(str(source))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password="dstu2026", owner_password="dstu-owner")
    with target.open("wb") as handle:
        writer.write(handle)


def make_blank_pdf(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=595, height=842)
    with path.open("wb") as handle:
        writer.write(handle)


def make_png_and_jpg(png_path: Path, jpg_path: Path) -> None:
    image = Image.new("RGB", (1600, 1000), "#EAF4F4")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((90, 90, 1510, 910), radius=48, fill="#FFFFFF", outline="#2B7A78", width=8)
    draw.text((150, 160), "知识地图：从问题到证据", fill="#17252A", font=font(58, True))
    nodes = [(300, 430, "提出问题"), (800, 340, "收集证据"), (1230, 500, "形成解释"), (760, 730, "反思修正")]
    for x, y, label in nodes:
        draw.ellipse((x - 145, y - 70, x + 145, y + 70), fill="#3AAFA9")
        box = draw.textbbox((0, 0), label, font=font(30, True))
        draw.text((x - (box[2] - box[0]) / 2, y - 20), label, fill="white", font=font(30, True))
    for start, end in zip(nodes, nodes[1:] + nodes[:1]):
        draw.line((start[0], start[1], end[0], end[1]), fill="#FE6D73", width=12)
    draw.text((150, 835), f"{MARKER}-IMAGE", fill="#4E5D5D", font=font(27))
    image.save(png_path)
    image.save(jpg_path, quality=88, optimize=True)


def make_xlsx(path: Path) -> None:
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
<Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
</Types>"""
    root_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>"""
    workbook = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<sheets><sheet name="学习计划" sheetId="1" r:id="rId1"/></sheets>
</workbook>"""
    workbook_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>"""
    rows = [
        ["主题", "计划分钟", "完成度", "校验标识"],
        ["数学", "45", "80%", MARKER],
        ["语言", "30", "65%", "中文/English"],
        ["科学", "60", "90%", "H2O / CO2"],
    ]
    xml_rows = []
    for row_index, row in enumerate(rows, 1):
        cells = "".join(
            f'<c r="{chr(65 + col)}{row_index}" t="inlineStr"><is><t>{value}</t></is></c>'
            for col, value in enumerate(row)
        )
        xml_rows.append(f'<row r="{row_index}">{cells}</row>')
    sheet = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        f"<sheetData>{''.join(xml_rows)}</sheetData></worksheet>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", root_rels)
        archive.writestr("xl/workbook.xml", workbook)
        archive.writestr("xl/_rels/workbook.xml.rels", workbook_rels)
        archive.writestr("xl/worksheets/sheet1.xml", sheet)


def make_pptx(path: Path) -> None:
    deck = Presentation()
    deck.slide_width = PptxInches(13.333)
    deck.slide_height = PptxInches(7.5)
    colors = [("173F5F", "F6D55C"), ("20639B", "FFFFFF"), ("3CAEA3", "17252A")]
    content = [
        ("学习资源预览测试", f"演示文稿样本｜{MARKER}"),
        ("三步学习法", "1 设定问题\n2 主动检索\n3 间隔复习"),
        ("完成检查", "文字清晰｜页面可切换｜比例正确"),
    ]
    for (title, body), (background, accent) in zip(content, colors):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = PptxRGBColor.from_string(background)
        shape = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.9), PptxInches(11.8), PptxInches(1.2))
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = PptxPt(38)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PptxRGBColor.from_string(accent)
        body_shape = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(2.5), PptxInches(10.8), PptxInches(3.3))
        body_frame = body_shape.text_frame
        body_frame.word_wrap = True
        for index, line in enumerate(body.split("\n")):
            p = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            p.text = line
            p.font.size = PptxPt(24)
            p.font.color.rgb = PptxRGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT
    deck.save(path)


def make_epub(path: Path) -> None:
    container = """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
<rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles>
</container>"""
    opf = f"""<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="bookid">
<metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:identifier id="bookid">dstu-fixture</dc:identifier><dc:title>学习资源 EPUB 样本</dc:title><dc:language>zh-CN</dc:language></metadata>
<manifest><item id="chapter" href="chapter.xhtml" media-type="application/xhtml+xml"/></manifest>
<spine><itemref idref="chapter"/></spine></package>"""
    chapter = f"""<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml"><head><title>样本章节</title></head>
<body><h1>第一章：主动学习</h1><p>{MARKER}-EPUB</p><p>这是用于上传与预览测试的最小 EPUB。</p></body></html>"""
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        archive.writestr("META-INF/container.xml", container)
        archive.writestr("OEBPS/content.opf", opf)
        archive.writestr("OEBPS/chapter.xhtml", chapter)


def make_wav(path: Path) -> None:
    sample_rate = 22050
    duration = 1.5
    with wave.open(str(path), "wb") as audio:
        audio.setnchannels(1)
        audio.setsampwidth(2)
        audio.setframerate(sample_rate)
        frames = bytearray()
        for i in range(int(sample_rate * duration)):
            value = int(9000 * math.sin(2 * math.pi * 440 * i / sample_rate))
            frames.extend(struct.pack("<h", value))
        audio.writeframes(frames)


def add_manifest_entry(entries: list[dict[str, object]], path: Path, category: str, valid: bool, expectation: str) -> None:
    entries.append(
        {
            "file": path.name,
            "category": category,
            "valid": valid,
            "bytes": path.stat().st_size,
            "expectation": expectation,
        }
    )


def main() -> None:
    if OUT.exists():
        shutil.rmtree(OUT)
    TMP.mkdir(parents=True)
    entries: list[dict[str, object]] = []

    illustration = TMP / "docx-illustration.png"
    make_illustration(illustration)

    rich_docx = OUT / "01-富文本多页-气候学习笔记.docx"
    table_docx = OUT / "02-横向长表格-45行学习计划.docx"
    media_docx = OUT / "03-图文与多语言-特殊字符.docx"
    empty_docx = OUT / "04-空白文档.docx"
    corrupt_docx = OUT / "05-损坏文档.docx"
    make_rich_docx(rich_docx, illustration)
    make_table_docx(table_docx)
    make_media_docx(media_docx, illustration)
    make_empty_docx(empty_docx)
    corrupt_docx.write_bytes(b"not-a-valid-docx\x00\x01")
    for path, valid, expectation in [
        (rich_docx, True, "上传成功；可分页预览文本、表格和图片"),
        (table_docx, True, "上传成功；横向宽表不裁切，可滚动/缩放"),
        (media_docx, True, "上传成功；图片、多语言和特殊字符可见"),
        (empty_docx, True, "上传成功；显示空白或无可提取内容提示"),
        (corrupt_docx, False, "明确拒绝或显示无法解析，不应卡死"),
    ]:
        add_manifest_entry(entries, path, "DOCX", valid, expectation)

    text_pdf = OUT / "06-文本型多页-可选择文字.pdf"
    table_pdf = OUT / "07-横向宽表格.pdf"
    scan_pdf = OUT / "08-扫描型双页-无文本层.pdf"
    encrypted_pdf = OUT / "09-密码保护-密码dstu2026.pdf"
    blank_pdf = OUT / "10-单页空白.pdf"
    corrupt_pdf = OUT / "11-损坏文件.pdf"
    make_text_pdf(text_pdf)
    make_table_pdf(table_pdf)
    make_scanned_pdf(scan_pdf)
    make_encrypted_pdf(text_pdf, encrypted_pdf)
    make_blank_pdf(blank_pdf)
    corrupt_pdf.write_bytes(b"%PDF-1.7\ncorrupt and truncated")
    for path, valid, expectation in [
        (text_pdf, True, "上传成功；三页可切换，中文文本可见"),
        (table_pdf, True, "上传成功；横向页面适配容器"),
        (scan_pdf, True, "上传成功；双页图像清晰可缩放"),
        (encrypted_pdf, True, "要求密码或明确提示受保护，不应白屏"),
        (blank_pdf, True, "上传成功；显示空白页"),
        (corrupt_pdf, False, "明确拒绝或显示 PDF 损坏，不应卡死"),
    ]:
        add_manifest_entry(entries, path, "PDF", valid, expectation)

    png = OUT / "12-知识地图-1600x1000.png"
    jpg = OUT / "13-知识地图-压缩照片.jpg"
    make_png_and_jpg(png, jpg)
    svg = OUT / "14-矢量流程图.svg"
    svg.write_text(
        f"""<svg xmlns="http://www.w3.org/2000/svg" width="900" height="500" viewBox="0 0 900 500">
<rect width="900" height="500" fill="#f5f3eb"/><circle cx="180" cy="250" r="90" fill="#3caea3"/>
<rect x="360" y="160" width="180" height="180" rx="28" fill="#20639b"/>
<path d="M650 340 L740 160 L830 340 Z" fill="#ed553b"/>
<text x="60" y="70" font-size="32" font-family="sans-serif" fill="#173f5f">SVG 学习流程｜{MARKER}</text>
</svg>""",
        encoding="utf-8",
    )
    for path in [png, jpg, svg]:
        add_manifest_entry(entries, path, "IMAGE", True, "上传成功；图像比例正确且可缩放")

    markdown = OUT / "15-Markdown-公式代码与表格.md"
    markdown.write_text(
        f"""# 学习资源 Markdown 样本

校验标识：`{MARKER}-MD`

## 内容覆盖

- 中文与 **粗体**
- [链接](https://example.com)
- 行内公式：$E = mc^2$

| 科目 | 分钟 | 状态 |
| --- | ---: | --- |
| 数学 | 45 | 完成 |
| 英语 | 30 | 进行中 |

```python
def spaced_repetition(days: int) -> str:
    return f"next review in {{days}} days"
```
""",
        encoding="utf-8",
    )
    text = OUT / "16-UTF8纯文本-中英混排.txt"
    text.write_text(f"{MARKER}-TXT\n第一行：简体中文。\nSecond line: English.\n第三行：制表符\t和特殊字符 ✓。\n", encoding="utf-8")
    long_name = OUT / ("17-" + "超长文件名-" * 12 + "边界测试.txt")
    long_name.write_text(f"{MARKER}-LONG-NAME\n长文件名不应破坏列表布局。", encoding="utf-8")
    html = OUT / "18-HTML-语义结构.html"
    html.write_text(
        f"<!doctype html><html lang='zh-CN'><meta charset='utf-8'><title>测试</title><body><h1>HTML 学习资料</h1><p>{MARKER}-HTML</p><details><summary>展开答案</summary><p>主动回忆。</p></details></body></html>",
        encoding="utf-8",
    )
    data_json = OUT / "19-JSON-学习进度.json"
    data_json.write_text(
        json.dumps({"marker": MARKER, "subject": "数学", "progress": 0.75, "tags": ["复习", "函数"]}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    data_csv = OUT / "20-CSV-学习计划.csv"
    with data_csv.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerows([["主题", "时长", "状态"], ["代数", 45, "完成"], ["阅读", 30, "进行中"], ["标识", MARKER, ""]])
    for path, category in [(markdown, "MARKDOWN"), (text, "TEXT"), (long_name, "TEXT"), (html, "HTML"), (data_json, "JSON"), (data_csv, "CSV")]:
        add_manifest_entry(entries, path, category, True, "上传成功；文本编码正确，长内容可滚动")

    xlsx = OUT / "21-XLSX-学习计划.xlsx"
    pptx = OUT / "22-PPTX-三页学习法.pptx"
    epub = OUT / "23-EPUB-最小电子书.epub"
    wav = OUT / "24-WAV-440Hz音频.wav"
    archive_path = OUT / "25-ZIP-资料包.zip"
    make_xlsx(xlsx)
    make_pptx(pptx)
    make_epub(epub)
    make_wav(wav)
    with zipfile.ZipFile(archive_path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("README.txt", f"{MARKER}-ZIP\n压缩包内部文本。")
        archive.writestr("notes/chapter-1.md", "# 第一章\n主动回忆与间隔复习。")
    for path, category, expectation in [
        (xlsx, "XLSX", "如支持则显示工作表；否则明确标记暂不支持"),
        (pptx, "PPTX", "如支持则显示三页幻灯片；否则明确标记暂不支持"),
        (epub, "EPUB", "如支持则显示章节；否则明确标记暂不支持"),
        (wav, "AUDIO", "如支持则显示播放器；否则明确标记暂不支持"),
        (archive_path, "ARCHIVE", "允许作为资源保存或明确拒绝，不应误当文本"),
    ]:
        add_manifest_entry(entries, path, category, True, expectation)

    manifest = {
        "suite": "Deep Student learning resource upload/preview fixtures",
        "marker": MARKER,
        "passwords": {"09-密码保护-密码dstu2026.pdf": "dstu2026"},
        "count": len(entries),
        "files": entries,
    }
    (OUT / "00-测试清单-manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    shutil.rmtree(TMP)
    print(f"Generated {len(entries)} fixtures in {OUT}")


if __name__ == "__main__":
    main()
